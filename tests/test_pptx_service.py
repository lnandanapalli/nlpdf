"""Tests for PowerPoint (PPTX) to PDF conversion service."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.services.pptx_service import pptx_to_pdf


@pytest.fixture()
def pptx_file(tmp_path: Path) -> Path:
    """Create a sample PPTX with text slides."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "This is the first slide."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "More content here."

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def pptx_with_table(tmp_path: Path) -> Path:
    """Create a PPTX with a table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    rows, cols = 3, 2
    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(6), Inches(2)
    ).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "B"
    table.cell(2, 1).text = "200"

    path = tmp_path / "table.pptx"
    prs.save(str(path))
    return path


class TestPptxToPdf:
    """Tests for pptx_to_pdf."""

    def test_basic_conversion(self, pptx_file: Path, tmp_path: Path):
        output = tmp_path / "output.pdf"
        result = pptx_to_pdf(pptx_file, output)

        assert result == output
        assert output.exists()
        assert output.stat().st_size > 0
        assert output.read_bytes()[:5] == b"%PDF-"

    def test_table_conversion(self, pptx_with_table: Path, tmp_path: Path):
        output = tmp_path / "output.pdf"
        result = pptx_to_pdf(pptx_with_table, output)

        assert result == output
        assert output.exists()
        assert output.read_bytes()[:5] == b"%PDF-"

    def test_letter_paper_size(self, pptx_file: Path, tmp_path: Path):
        output = tmp_path / "output.pdf"
        result = pptx_to_pdf(pptx_file, output, paper_size="letter")

        assert result == output
        assert output.exists()
        assert output.read_bytes()[:5] == b"%PDF-"

    def test_empty_pptx(self, tmp_path: Path):
        prs = Presentation()
        path = tmp_path / "empty.pptx"
        prs.save(str(path))
        output = tmp_path / "output.pdf"

        result = pptx_to_pdf(path, output)
        assert result == output
        assert output.exists()
