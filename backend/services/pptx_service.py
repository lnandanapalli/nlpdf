"""PowerPoint (PPTX) to PDF conversion service using python-pptx and xhtml2pdf."""

import structlog
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from xhtml2pdf import pisa

from backend.services.pdf_fonts import register_fonts

logger = structlog.get_logger(__name__)

CSS = """\
body {
    font-family: DejaVuSans, Helvetica, Arial, sans-serif;
    font-size: 14px;
    line-height: 1.5;
    color: #222;
    margin: 0;
    padding: 0;
}
.slide {
    padding: 40px 50px;
    page-break-after: always;
}
.slide:last-child {
    page-break-after: avoid;
}
.slide-title {
    font-size: 28px;
    font-weight: bold;
    margin-bottom: 20px;
    color: #1a1a2e;
}
.slide-content {
    font-size: 16px;
    margin-bottom: 8px;
}
table {
    border-collapse: collapse;
    width: 100%;
    margin: 12px 0;
}
th, td {
    border: 1px solid #ddd;
    padding: 8px;
    text-align: left;
}
th {
    background-color: #f4f4f4;
    font-weight: bold;
}
.slide-number {
    font-size: 10px;
    color: #999;
    text-align: right;
    margin-top: 20px;
}
"""

PAPER_SIZES = {
    "A4": "@page { size: A4 landscape; margin: 1.5cm; }",
    "letter": "@page { size: letter landscape; margin: 0.75in; }",
}


def _extract_table_html(table) -> str:
    """Convert a PPTX table shape to an HTML table."""
    rows = []
    for i, row in enumerate(table.rows):
        tag = "th" if i == 0 else "td"
        cells = "".join(f"<{tag}>{cell.text}</{tag}>" for cell in row.cells)
        rows.append(f"<tr>{cells}</tr>")
    return f"<table>{''.join(rows)}</table>"


def _extract_slide_html(slide, slide_num: int) -> str:
    """Extract content from a single slide as HTML."""
    parts: list[str] = []

    for shape in slide.shapes:
        if shape.has_table:
            parts.append(_extract_table_html(shape.table))
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                # Heuristic: large font or first text block = title
                is_title = False
                if paragraph.runs:
                    font = paragraph.runs[0].font
                    if font.size and font.size > Emu(200000):  # ~14pt+
                        is_title = True
                    if font.bold:
                        is_title = True

                if is_title and not any("slide-title" in p for p in parts):
                    parts.append(f'<div class="slide-title">{text}</div>')
                else:
                    parts.append(f'<div class="slide-content">{text}</div>')

    content = "\n".join(parts) if parts else '<div class="slide-content">&nbsp;</div>'

    return (
        f'<div class="slide">'
        f"{content}"
        f'<div class="slide-number">Slide {slide_num}</div>'
        f"</div>"
    )


def pptx_to_pdf(input_path: Path, output_path: Path, paper_size: str = "A4") -> Path:
    """
    Convert a PowerPoint (PPTX) file to PDF.

    Extracts text and table content from each slide and renders
    them as styled HTML pages via xhtml2pdf.

    Args:
        input_path: Path to the input PPTX file.
        output_path: Path to save the output PDF.
        paper_size: Paper size — "A4" or "letter".

    Returns:
        Path to the generated PDF file.

    Raises:
        ValueError: If the PPTX file cannot be read or converted.
    """
    register_fonts()

    prs = Presentation(str(input_path))

    slides_html = []
    for i, slide in enumerate(prs.slides, start=1):
        slides_html.append(_extract_slide_html(slide, i))

    page_css = PAPER_SIZES.get(paper_size, PAPER_SIZES["A4"])
    body = "\n".join(slides_html)

    html = (
        "<!DOCTYPE html>"
        "<html><head><meta charset='utf-8'/>"
        f"<style>{page_css}\n{CSS}</style>"
        f"</head><body>{body}</body></html>"
    )

    with open(output_path, "wb") as f:
        status = pisa.CreatePDF(html, dest=f)

    if status.err:
        output_path.unlink(missing_ok=True)
        raise ValueError(f"PDF conversion failed with {status.err} error(s)")

    logger.info(
        "PowerPoint converted to PDF",
        input=str(input_path),
        output=str(output_path),
        slides=len(slides_html),
        paper_size=paper_size,
    )

    return output_path
